import os
import argparse
from typing import Dict, Any, List

from dotenv import load_dotenv
load_dotenv()

from .chunk import chunk_text
from .embed import Embedder
from .store import get_store

def _read_text_from_file(path: str) -> List[Dict[str, Any]]:
    ext = os.path.splitext(path)[1].lower()
    records = []
    if ext in [".md", ".txt"]:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            txt = f.read()
        records.append({"text": txt, "metadata": {"doc_title": os.path.basename(path), "source_path": path}})
    elif ext == ".pdf":
        try:
            import fitz  # pymupdf
            doc = fitz.open(path)
            for i, page in enumerate(doc):
                txt = page.get_text()
                records.append({"text": txt, "metadata": {"doc_title": os.path.basename(path), "source_path": path, "page": i+1}})
        except Exception as e:
            print(f"[WARN] PDF read failed for {path}: {e}")
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(path)
            for i, slide in enumerate(prs.slides):
                txt = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        txt.append(shape.text)
                slide_txt = "\n".join(txt)
                records.append({"text": slide_txt, "metadata": {"doc_title": os.path.basename(path), "source_path": path, "page": i+1, "slide_title": f"Slide {i+1}"}})
        except Exception as e:
            print(f"[WARN] PPTX read failed for {path}: {e}")
    else:
        print(f"[SKIP] Unsupported file type: {path}")
    return records

def ingest_path(path: str, chunk_size: int = 900, chunk_overlap: int = 120):
    store = get_store()
    embedder = Embedder()
    to_upsert = {"ids": [], "documents": [], "embeddings": [], "metadatas": []}
    did = 0

    for root, _, files in os.walk(path):
        for name in files:
            fpath = os.path.join(root, name)
            recs = _read_text_from_file(fpath)
            for rec in recs:
                for ch_i, ch in enumerate(chunk_text(rec["text"], chunk_size=chunk_size, chunk_overlap=chunk_overlap)):
                    _id = f"{name}_{did}_{ch_i}"
                    did += 1
                    to_upsert["ids"].append(_id)
                    to_upsert["documents"].append(ch)
                    to_upsert["metadatas"].append(rec["metadata"])

    if not to_upsert["documents"]:
        print("[INFO] No documents to embed.")
        return

    # Batch embed to avoid token/time limits
    batch = 64
    embeddings = []
    docs = to_upsert["documents"]
    for i in range(0, len(docs), batch):
        chunk = docs[i:i+batch]
        embs = embedder.embed(chunk)
        embeddings.extend(embs)
    to_upsert["embeddings"] = embeddings

    store.add(
        ids=to_upsert["ids"],
        documents=to_upsert["documents"],
        metadatas=to_upsert["metadatas"],
        embeddings=to_upsert["embeddings"],
    )
    print(f"[OK] Ingested {len(to_upsert['documents'])} chunks.")
    return len(to_upsert["documents"])

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--path", type=str, required=True, help="Folder containing PDFs/PPTX/MD/TXT")
    ap.add_argument("--chunk_size", type=int, default=900)
    ap.add_argument("--chunk_overlap", type=int, default=120)
    args = ap.parse_args()
    ingest_path(args.path, args.chunk_size, args.chunk_overlap)
